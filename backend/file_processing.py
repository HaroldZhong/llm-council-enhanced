"""
File processing logic for the LLM Council.

Handles text extraction from various file types:
- PDF, DOCX, PPTX, XLSX, CSV, TXT, MD, HTML, JSON
- Images via vision models

Returns structured extraction results with confidence metrics.
"""

import io
import csv
import json
from typing import Dict, Any, Optional, List
from fastapi import UploadFile
import pypdf
import docx

from .openrouter import query_model
from .config import OPENROUTER_API_KEY
from .logger import logger

# Limits
MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB (increased for office docs)
MAX_FILE_SIZE_MB = MAX_FILE_SIZE / (1024 * 1024)
MAX_TEXT_LENGTH = 100000  # 100k characters (increased for multi-page docs)

# Confidence thresholds for PDF extraction
MIN_CHARS_THRESHOLD = 500  # If less than this, likely a scanned PDF
MAX_EMPTY_PAGE_RATIO = 0.5  # If more than half pages empty, needs OCR


class ExtractionResult:
    """Structured extraction result."""
    
    def __init__(
        self,
        status: str = "success",
        text: str = "",
        method: str = "local",
        warning: Optional[str] = None,
        error: Optional[str] = None,
        stats: Optional[Dict[str, Any]] = None
    ):
        self.status = status
        self.text = text
        self.method = method
        self.warning = warning
        self.error = error
        self.stats = stats or {}
    
    def to_dict(self) -> Dict[str, Any]:
        return {
            "status": self.status,
            "text": self.text,
            "method": self.method,
            "warning": self.warning,
            "error": self.error,
            "stats": self.stats
        }


def get_mime_type(filename: str, content_type: Optional[str]) -> str:
    """Normalize MIME type from filename extension if content_type is generic."""
    if content_type and content_type != "application/octet-stream":
        return content_type
    
    ext = filename.lower().split('.')[-1] if '.' in filename else ''
    
    mime_map = {
        'pdf': 'application/pdf',
        'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        'xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        'csv': 'text/csv',
        'txt': 'text/plain',
        'md': 'text/markdown',
        'html': 'text/html',
        'htm': 'text/html',
        'json': 'application/json',
        'jpg': 'image/jpeg',
        'jpeg': 'image/jpeg',
        'png': 'image/png',
        'gif': 'image/gif',
        'webp': 'image/webp',
    }
    
    return mime_map.get(ext, content_type or 'application/octet-stream')


async def process_file(content: bytes, filename: str, mime_type: str) -> ExtractionResult:
    """
    Process file content and extract text.
    
    Args:
        content: Raw file bytes
        filename: Original filename
        mime_type: MIME type of file
    
    Returns:
        ExtractionResult with status, text, and metadata
    """
    # Normalize MIME type
    mime_type = get_mime_type(filename, mime_type)
    
    logger.info(f"[EXTRACT] Processing {filename} ({mime_type}, {len(content)} bytes)")
    
    # Size check
    if len(content) > MAX_FILE_SIZE:
        return ExtractionResult(
            status="failed",
            error=f"File too large ({len(content) / 1024 / 1024:.1f}MB). Max: {MAX_FILE_SIZE_MB:.0f}MB."
        )
    
    try:
        # Route to appropriate extractor
        if mime_type == 'application/pdf':
            return _extract_pdf(content)
        
        elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
            return _extract_docx(content)
        
        elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
            return _extract_pptx(content)
        
        elif mime_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
            return _extract_xlsx(content)
        
        elif mime_type == 'text/csv':
            return _extract_csv(content)
        
        elif mime_type in ('text/plain', 'text/markdown'):
            return _extract_text(content)
        
        elif mime_type == 'text/html':
            return _extract_html(content)
        
        elif mime_type == 'application/json':
            return _extract_json(content)
        
        elif mime_type.startswith('image/'):
            return await _extract_image(content, mime_type)
        
        else:
            return ExtractionResult(
                status="failed",
                error=f"Unsupported file type: {mime_type}"
            )
    
    except Exception as e:
        logger.error(f"[EXTRACT] Error processing {filename}: {e}")
        return ExtractionResult(
            status="failed",
            error=f"Extraction failed: {str(e)}"
        )


def _extract_pdf(content: bytes) -> ExtractionResult:
    """Extract text from PDF with confidence metrics."""
    pdf_file = io.BytesIO(content)
    reader = pypdf.PdfReader(pdf_file)
    
    page_texts = []
    empty_pages = 0
    
    for i, page in enumerate(reader.pages):
        try:
            text = page.extract_text() or ""
            page_texts.append(f"[Page {i + 1}]\n{text}" if text.strip() else "")
            if not text.strip():
                empty_pages += 1
        except Exception as e:
            logger.warning(f"[EXTRACT] Failed to extract page {i + 1}: {e}")
            page_texts.append("")
            empty_pages += 1
    
    full_text = "\n\n".join(filter(None, page_texts))
    page_count = len(reader.pages)
    char_count = len(full_text)
    empty_ratio = empty_pages / page_count if page_count > 0 else 1.0
    
    stats = {
        "page_count": page_count,
        "char_count": char_count,
        "empty_page_ratio": empty_ratio
    }
    
    # Determine status based on extraction quality
    if char_count < MIN_CHARS_THRESHOLD or empty_ratio > MAX_EMPTY_PAGE_RATIO:
        return ExtractionResult(
            status="partial",
            text=full_text,
            method="local",
            warning="This PDF may be scanned or image-based. Enhanced extraction (OCR) is recommended for better results.",
            stats=stats
        )
    
    # Truncate if needed
    if len(full_text) > MAX_TEXT_LENGTH:
        full_text = full_text[:MAX_TEXT_LENGTH] + "\n\n[...content truncated...]"
        return ExtractionResult(
            status="success",
            text=full_text,
            method="local",
            warning=f"Content truncated to {MAX_TEXT_LENGTH} characters.",
            stats=stats
        )
    
    return ExtractionResult(
        status="success",
        text=full_text,
        method="local",
        stats=stats
    )


def _extract_docx(content: bytes) -> ExtractionResult:
    """Extract text from DOCX."""
    docx_file = io.BytesIO(content)
    doc = docx.Document(docx_file)
    
    paragraphs = []
    for para in doc.paragraphs:
        if para.text.strip():
            paragraphs.append(para.text)
    
    # Also extract tables
    table_texts = []
    for table in doc.tables:
        table_rows = []
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                table_rows.append(row_text)
        if table_rows:
            table_texts.append("\n".join(table_rows))
    
    full_text = "\n\n".join(paragraphs)
    if table_texts:
        full_text += "\n\n[Tables]\n" + "\n\n".join(table_texts)
    
    return ExtractionResult(
        status="success",
        text=full_text,
        method="local",
        stats={"char_count": len(full_text)}
    )


def _extract_pptx(content: bytes) -> ExtractionResult:
    """Extract text from PPTX (PowerPoint)."""
    try:
        from pptx import Presentation
    except ImportError:
        return ExtractionResult(
            status="failed",
            error="PPTX support not installed. Run: pip install python-pptx"
        )
    
    pptx_file = io.BytesIO(content)
    prs = Presentation(pptx_file)
    
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        slide_parts = [f"[Slide {i}]"]
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_parts.append(shape.text.strip())
        
        # Extract speaker notes if present
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text
            if notes_text.strip():
                slide_parts.append(f"[Notes] {notes_text.strip()}")
        
        slides_text.append("\n".join(slide_parts))
    
    full_text = "\n\n".join(slides_text)
    
    return ExtractionResult(
        status="success",
        text=full_text,
        method="local",
        stats={"slide_count": len(prs.slides), "char_count": len(full_text)}
    )


def _extract_xlsx(content: bytes) -> ExtractionResult:
    """Extract data from XLSX (Excel)."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        return ExtractionResult(
            status="failed",
            error="XLSX support not installed. Run: pip install openpyxl"
        )
    
    xlsx_file = io.BytesIO(content)
    wb = load_workbook(xlsx_file, read_only=True, data_only=True)
    
    sheets_text = []
    sheet_count = 0
    
    for sheet_name in wb.sheetnames:
        sheet_count += 1
        ws = wb[sheet_name]
        
        rows = []
        row_count = 0
        for row in ws.iter_rows(max_row=1000):  # Limit rows to avoid huge outputs
            row_values = [str(cell.value) if cell.value is not None else "" for cell in row]
            if any(v.strip() for v in row_values):
                rows.append(" | ".join(row_values))
                row_count += 1
        
        if rows:
            sheet_header = f"[Sheet: {sheet_name}] ({row_count} rows)"
            sheets_text.append(f"{sheet_header}\n" + "\n".join(rows[:100]))  # Limit displayed rows
            if row_count > 100:
                sheets_text[-1] += f"\n[...{row_count - 100} more rows...]"
    
    wb.close()
    full_text = "\n\n".join(sheets_text)
    
    return ExtractionResult(
        status="success",
        text=full_text,
        method="local",
        stats={"sheet_count": sheet_count, "char_count": len(full_text)}
    )


def _extract_csv(content: bytes) -> ExtractionResult:
    """Extract data from CSV."""
    text = content.decode('utf-8', errors='replace')
    
    # Try to detect delimiter
    try:
        dialect = csv.Sniffer().sniff(text[:4096])
        delimiter = dialect.delimiter
    except csv.Error:
        delimiter = ','
    
    reader = csv.reader(io.StringIO(text), delimiter=delimiter)
    rows = []
    row_count = 0
    
    for row in reader:
        if row_count >= 1000:  # Limit rows
            break
        if any(cell.strip() for cell in row):
            rows.append(" | ".join(row))
            row_count += 1
    
    full_text = "\n".join(rows[:100])
    if row_count > 100:
        full_text += f"\n[...{row_count - 100} more rows...]"
    
    return ExtractionResult(
        status="success",
        text=full_text,
        method="local",
        stats={"row_count": row_count, "char_count": len(full_text)}
    )


def _extract_text(content: bytes) -> ExtractionResult:
    """Extract plain text / markdown."""
    text = content.decode('utf-8', errors='replace')
    
    if len(text) > MAX_TEXT_LENGTH:
        text = text[:MAX_TEXT_LENGTH] + "\n\n[...content truncated...]"
        return ExtractionResult(
            status="success",
            text=text,
            method="local",
            warning=f"Content truncated to {MAX_TEXT_LENGTH} characters.",
            stats={"char_count": len(text)}
        )
    
    return ExtractionResult(
        status="success",
        text=text,
        method="local",
        stats={"char_count": len(text)}
    )


def _extract_html(content: bytes) -> ExtractionResult:
    """Extract text from HTML."""
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        return ExtractionResult(
            status="failed",
            error="HTML parsing not installed. Run: pip install beautifulsoup4"
        )
    
    html = content.decode('utf-8', errors='replace')
    soup = BeautifulSoup(html, 'html.parser')
    
    # Remove script and style elements
    for element in soup(['script', 'style', 'nav', 'footer', 'header']):
        element.decompose()
    
    # Extract text
    text = soup.get_text(separator='\n', strip=True)
    
    if len(text) > MAX_TEXT_LENGTH:
        text = text[:MAX_TEXT_LENGTH] + "\n\n[...content truncated...]"
    
    return ExtractionResult(
        status="success",
        text=text,
        method="local",
        stats={"char_count": len(text)}
    )


def _extract_json(content: bytes) -> ExtractionResult:
    """Extract and format JSON."""
    text = content.decode('utf-8', errors='replace')
    
    try:
        data = json.loads(text)
        # Pretty print with truncation for large arrays
        formatted = json.dumps(data, indent=2, ensure_ascii=False)
        
        if len(formatted) > MAX_TEXT_LENGTH:
            formatted = formatted[:MAX_TEXT_LENGTH] + "\n\n[...JSON truncated...]"
        
        return ExtractionResult(
            status="success",
            text=formatted,
            method="local",
            stats={"char_count": len(formatted)}
        )
    except json.JSONDecodeError as e:
        return ExtractionResult(
            status="partial",
            text=text[:MAX_TEXT_LENGTH],
            method="local",
            warning=f"Invalid JSON: {e}. Returning raw content.",
            stats={"char_count": len(text)}
        )


async def _extract_image(content: bytes, mime_type: str) -> ExtractionResult:
    """Describe image using vision model."""
    import base64
    
    base64_image = base64.b64encode(content).decode("utf-8")
    data_url = f"data:{mime_type};base64,{base64_image}"
    
    messages = [
        {
            "role": "user",
            "content": [
                {
                    "type": "text",
                    "text": "Describe this image in detail. Focus on text, diagrams, charts, and key visual elements that would be relevant for analysis and discussion."
                },
                {
                    "type": "image_url",
                    "image_url": {"url": data_url}
                }
            ]
        }
    ]
    
    try:
        response = await query_model("google/gemini-2.5-flash", messages, timeout=30.0)
        
        if response and response.get("content"):
            description = response["content"]
            return ExtractionResult(
                status="success",
                text=f"[Image Description]\n{description}",
                method="vision",
                stats={"char_count": len(description)}
            )
    except Exception as e:
        logger.error(f"[EXTRACT] Vision model error: {e}")
    
    return ExtractionResult(
        status="failed",
        error="Failed to analyze image with vision model."
    )


# Legacy function for backwards compatibility
async def extract_text_from_file(file: UploadFile) -> Dict[str, Any]:
    """
    Legacy wrapper for backwards compatibility.
    Use process_file() for new code.
    """
    content = await file.read()
    mime_type = get_mime_type(file.filename, file.content_type)
    
    result = await process_file(content, file.filename, mime_type)
    
    # Convert to legacy format
    return {
        "text": result.text,
        "truncated": "truncated" in (result.warning or "").lower(),
        "error": result.error
    }
